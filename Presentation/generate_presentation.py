from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # Define custom colors
    DARK_BG = RGBColor(33, 33, 33)
    LIGHT_TEXT = RGBColor(240, 240, 240)
    ACCENT_COLOR = RGBColor(0, 120, 215) # Blue accent

    # Helper function to apply dark theme to a slide
    def apply_dark_theme(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    # Helper function to format text
    def format_text(shape, color=LIGHT_TEXT, size=None, bold=False):
        if not shape.has_text_frame:
            return
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = color
                if size:
                    run.font.size = size
                if bold:
                    run.font.bold = bold

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    apply_dark_theme(slide1)
    
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]

    title.text = "Applied AI & Robotics Portfolio"
    format_text(title, color=ACCENT_COLOR, bold=True)
    
    subtitle.text = "Javon Darby\nHouston Community College\nDeep Learning (ITAI 2376)\ngithub.com/JavonDarby/Javon-Darby-AI-Portfolio"
    format_text(subtitle, color=LIGHT_TEXT)

    # Slide 2: Professional Summary & Skills
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    apply_dark_theme(slide2)

    title2 = slide2.shapes.title
    title2.text = "Professional Profile & Technical Arsenal"
    format_text(title2, color=ACCENT_COLOR, bold=True)

    body_shape2 = slide2.placeholders[1]
    tf = body_shape2.text_frame
    tf.text = "Dedicated Applied AI and Robotics student focused on building scalable machine learning solutions from the ground up."
    
    p = tf.add_paragraph()
    p.text = "Core Competencies:"
    p.level = 0
    
    skills = [
        "Deep Neural Networks (PyTorch, TensorFlow, Keras)",
        "Generative AI & Probabilistic Diffusion Models",
        "Sequence Modeling (Transformers, LLMs, RNNs)",
        "Reinforcement Learning (Q-Learning, Gymnasium)"
    ]
    
    for skill in skills:
        p = tf.add_paragraph()
        p.text = skill
        p.level = 1
        
    format_text(body_shape2, color=LIGHT_TEXT)

    # Slide 3: Featured Project 1
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    apply_dark_theme(slide3)

    title3 = slide3.shapes.title
    title3.text = "Featured Project: Generative Diffusion Model"
    format_text(title3, color=ACCENT_COLOR, bold=True)

    body_shape3 = slide3.placeholders[1]
    tf3 = body_shape3.text_frame
    tf3.text = "Built a Denoising Diffusion Probabilistic Model (DDPM) entirely from scratch."

    points3 = [
        "Engineered a deep Convolutional U-Net architecture in PyTorch.",
        "Implemented forward (noise injection) and reverse (denoising) mathematical processes.",
        "Integrated Sinusoidal Positional Embeddings for timestep conditioning.",
        "Result: Successfully generated clear, high-fidelity digits from pure Gaussian noise, demonstrating mastery of modern generative AI concepts."
    ]

    for pt in points3:
        p = tf3.add_paragraph()
        p.text = pt
        p.level = 1
        
    format_text(body_shape3, color=LIGHT_TEXT)

    # Slide 4: Featured Project 2
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    apply_dark_theme(slide4)

    title4 = slide4.shapes.title
    title4.text = "Featured Project: Autonomous RL Agent"
    format_text(title4, color=ACCENT_COLOR, bold=True)

    body_shape4 = slide4.placeholders[1]
    tf4 = body_shape4.text_frame
    tf4.text = "Trained a Reinforcement Learning agent to solve the CartPole physics environment."

    points4 = [
        "Designed and implemented a Q-Learning algorithm from fundamental Bellman equations.",
        "Utilized Gymnasium for environment simulation and state-space discretization.",
        "Engineered an Epsilon-Greedy policy to balance exploration and exploitation.",
        "Result: Agent transitioned from random failures to a converged optimal policy, surviving the maximum 500 steps per episode."
    ]

    for pt in points4:
        p = tf4.add_paragraph()
        p.text = pt
        p.level = 1
        
    format_text(body_shape4, color=LIGHT_TEXT)

    # Slide 5: Closing
    slide5 = prs.slides.add_slide(title_slide_layout)
    apply_dark_theme(slide5)

    title5 = slide5.shapes.title
    subtitle5 = slide5.placeholders[1]

    title5.text = "Thank You"
    format_text(title5, color=ACCENT_COLOR, bold=True)
    
    subtitle5.text = "Please visit my GitHub to review the complete source code and documentation.\n\ngithub.com/JavonDarby/Javon-Darby-AI-Portfolio\n\nJavon Darby | Houston Community College"
    format_text(subtitle5, color=LIGHT_TEXT)

    # Save presentation
    output_path = "Pf_JavonDarby_ITAI2376.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
